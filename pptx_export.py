"""
Generator PPT dari data dashboard — mengikuti template presentasi perusahaan
(Madinah Group Indonesia / MFlash).

Dipakai oleh app.py: tombol "Buat PPT" memanggil build_deck(...) yang
mengembalikan file .pptx dalam bentuk bytes, siap di-download.
"""
import calendar
import io
from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

ASSETS = Path(__file__).parent / "assets"

# ---- palet mengikuti template perusahaan ----
NAVY = RGBColor(0x1F, 0x38, 0x64)      # judul & teks utama
NAVY_L = RGBColor(0x2E, 0x53, 0x94)
BLUE = RGBColor(0x2E, 0x9B, 0xD6)      # biru Madinah
ORANGE = RGBColor(0xF2, 0x8C, 0x1C)    # oranye MFlash
GREEN = RGBColor(0x4E, 0x7B, 0x3A)     # hijau bullet template
GREEN_B = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xC0, 0x39, 0x2B)
INK = RGBColor(0x20, 0x24, 0x2E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xDD, 0xE2, 0xEC)
CARD = RGBColor(0xF6, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG = RGBColor(0xF2, 0xF4, 0xF8)

F_H = "Calibri"
F_B = "Calibri"

BULAN_NAMES = ['', 'Januari', 'Februari', 'Maret', 'April', 'Mei', 'Juni', 'Juli',
               'Agustus', 'September', 'Oktober', 'November', 'Desember']
BULAN_SINGKAT = ['', 'Jan', 'Feb', 'Mar', 'Apr', 'Mei', 'Jun', 'Jul',
                 'Agu', 'Sep', 'Okt', 'Nov', 'Des']
HARI_NAMES = ['Senin', 'Selasa', 'Rabu', 'Kamis', 'Jumat', 'Sabtu', 'Minggu']


def tgl_id(d, *, pendek=False):
    """Format tanggal ke bahasa Indonesia, mis. '8 Maret 2026' / '8 Mar'."""
    if pd.isna(d):
        return "-"
    d = pd.Timestamp(d)
    if pendek:
        return f"{d.day} {BULAN_SINGKAT[d.month]}"
    return f"{d.day} {BULAN_NAMES[d.month]} {d.year}"


def nf(n):
    try:
        return f"{int(n):,}".replace(",", ".")
    except Exception:
        return str(n)


def dec(x, n=1):
    return f"{x:.{n}f}".replace(".", ",")


# =====================================================================
# primitives
# =====================================================================
def _bg_image(slide, prs):
    """Background bertekstur seperti template; fallback ke abu-abu polos."""
    p = ASSETS / "bg.jpg"
    if p.exists():
        slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
    else:
        f = slide.background.fill
        f.solid()
        f.fore_color.rgb = WHITE


def _text(slide, x, y, w, h, text, *, font=F_B, size=12, bold=False, color=INK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, line_sp=None):
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_sp:
            p.line_spacing = Pt(line_sp)
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if spacing:
            r.font._rPr.set("spc", str(int(spacing * 100)))
    return tb


def _rect(slide, x, y, w, h, fill, *, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          line_color=None, radius=0.06):
    sp = slide.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def _title(slide, text):
    """Judul bergaya template: segitiga biru + teks kapital navy."""
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, In(0.62), In(0.42), In(0.3), In(0.36))
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = NAVY
    tri.line.fill.background()
    tri.shadow.inherit = False
    _text(slide, 1.06, 0.36, 11.6, 0.52, text.upper(), font=F_H, size=28, bold=True, color=NAVY)


def _bullet(slide, x, y, w, text, *, size=13, color=GREEN, bold=True):
    """Bullet bulat kecil + teks, seperti bullet hijau di template."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y + 0.07), In(0.1), In(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    dot.shadow.inherit = False
    return _text(slide, x + 0.26, y, w - 0.26, 0.6, text, size=size, bold=bold, color=color)


def _footer(slide, txt):
    _text(slide, 0.62, 7.02, 12.1, 0.26, txt, size=9, color=MUTED)


def _logos(slide, *, small=True):
    """Logo perusahaan di pojok kanan atas (mengikuti template)."""
    m = ASSETS / "logo_madinah.png"
    f = ASSETS / "logo_mflash.png"
    if small:
        if f.exists():
            slide.shapes.add_picture(str(f), In(12.42), In(0.26), height=In(0.5))
        if m.exists():
            slide.shapes.add_picture(str(m), In(11.4), In(0.32), height=In(0.38))


def _style_chart(chart, *, colors=None, title=None, show_legend=False,
                 legend_pos=XL_LEGEND_POSITION.TOP, number_format=None,
                 show_labels=True, label_pos=XL_LABEL_POSITION.OUTSIDE_END,
                 val_max=None, cat_size=10, label_size=9):
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        r = chart.chart_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.name = F_B
        r.font.color.rgb = NAVY

    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
        chart.legend.font.name = F_B
        chart.legend.font.color.rgb = INK

    plot = chart.plots[0]
    plot.has_data_labels = show_labels
    if show_labels:
        dl = plot.data_labels
        dl.font.size = Pt(label_size)
        dl.font.name = F_B
        dl.font.color.rgb = INK
        if number_format:
            dl.number_format = number_format
            dl.number_format_is_linked = False
        if label_pos is not None:
            try:
                dl.position = label_pos
            except Exception:
                pass

    if colors:
        for i, ser in enumerate(chart.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = colors[i % len(colors)]

    try:
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(cat_size)
        ca.tick_labels.font.name = F_B
        ca.tick_labels.font.color.rgb = INK
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = LINE
    except Exception:
        pass
    try:
        va = chart.value_axis
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.name = F_B
        va.tick_labels.font.color.rgb = MUTED
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = LINE
        va.major_gridlines.format.line.width = Pt(0.75)
        va.format.line.fill.background()
        if val_max is not None:
            va.maximum_scale = val_max
        va.minimum_scale = 0
    except Exception:
        pass


def _fit_size(text, box_w_in, *, base=25, min_size=11):
    """Perkirakan ukuran font agar teks muat satu baris di dalam kartu.

    Lebar rata-rata karakter Calibri bold ≈ 0.52 × ukuran font (dalam poin).
    """
    t = str(text)
    if not t:
        return base
    usable_pt = box_w_in * 72.0
    size = min(base, usable_pt / (len(t) * 0.52))
    return max(min_size, round(size, 1))


def _kpi_row(slide, cards, y=1.4, h=1.3):
    """Deret kartu KPI berwarna."""
    n = len(cards)
    gap = 0.16
    total_w = 12.1
    w = (total_w - gap * (n - 1)) / n
    for i, c in enumerate(cards):
        x = 0.62 + i * (w + gap)
        _rect(slide, x, y, w, h, c.get("fill", CARD),
              line_color=None if c.get("fill") else LINE)
        tcol = c.get("text", WHITE if c.get("fill") else INK)
        sub_col = c.get("subcolor", tcol)
        _text(slide, x + 0.18, y + 0.13, w - 0.36, 0.26, c["label"],
              size=9, bold=True, color=sub_col, spacing=0.8)
        vsize = _fit_size(c["value"], w - 0.36)
        _text(slide, x + 0.18, y + 0.4, w - 0.36, 0.52, c["value"],
              font=F_H, size=vsize, bold=True, color=tcol, anchor=MSO_ANCHOR.MIDDLE)
        if c.get("sub"):
            _text(slide, x + 0.18, y + 0.96, w - 0.36, 0.28, c["sub"], size=9.5, color=sub_col)


# =====================================================================
# perhitungan
# =====================================================================
def _period_days(df):
    sub = df.dropna(subset=["TAHUN", "BULAN"])
    if sub.empty:
        return 0
    today = date.today()
    tot = 0
    for _, r in sub[["TAHUN", "BULAN"]].drop_duplicates().iterrows():
        y, m = int(r["TAHUN"]), int(r["BULAN"])
        d = calendar.monthrange(y, m)[1]
        if y == today.year and m == today.month:
            d = min(d, today.day)
        elif (y > today.year) or (y == today.year and m > today.month):
            d = 0
        tot += d
    return tot


def _jenis_pending(s):
    s = str(s).upper() if pd.notna(s) else ""
    if "TEKNISI" in s:
        return "Teknisi"
    if "CUSTOMER" in s:
        return "Customer"
    if "SPAREPART" in s:
        return "Sparepart"
    if "KLAIM" in s or "GARANSI" in s:
        return "Klaim Garansi"
    if "KOMPLAIN" in s or "COMPLAIN" in s:
        return "Komplain"
    return "Umum"


def rp(v, singkat=True):
    """Format rupiah ringkas: 1,2 M / 340,5 jt / 12.500."""
    try:
        v = float(v)
    except (TypeError, ValueError):
        return "-"
    neg, v = v < 0, abs(float(v))
    if singkat:
        if v >= 1_000_000_000:
            s = f"Rp {v/1_000_000_000:,.2f} M"
        elif v >= 1_000_000:
            s = f"Rp {v/1_000_000:,.1f} jt"
        elif v >= 1_000:
            s = f"Rp {v/1_000:,.0f} rb"
        else:
            s = f"Rp {v:,.0f}"
    else:
        s = f"Rp {v:,.0f}"
    s = s.replace(",", "#").replace(".", ",").replace("#", ".")
    return ("-" + s) if neg else s


KATA_KUNCI_TARIF = ['INTERFACE', 'NORMAL', 'MATI TOTAL', 'PROMO']
TARIF_AWAL_PPT = {'Interface': 20.0, 'Normal': 30.0, 'Mati Total': 32.0, 'Promo': 60.0}


def _label_tarif(nama_barang, prioritas='NORMAL'):
    s = str(nama_barang).upper()
    cocok = [k for k in KATA_KUNCI_TARIF if k in s]
    if not cocok:
        return 'Lainnya'
    urutan = [prioritas] + [k for k in KATA_KUNCI_TARIF if k != prioritas]
    for k in urutan:
        if k in cocok:
            return k.title()
    return cocok[0].title()


def _jenis_cancel(s):
    s = str(s).upper() if pd.notna(s) else ""
    if "DIAMBIL" in s:
        return "Tidak jadi diambil"
    if "TEKNISI" in s:
        return "Oleh teknisi"
    if "CUSTOMER" in s or "USER" in s:
        return "Dibatalkan customer"
    if "SPAREPART" in s:
        return "Sparepart"
    if "ADMIN" in s:
        return "Admin"
    return "Lainnya"


# =====================================================================
# SLIDE BUILDERS
# =====================================================================
def _slide_cover(prs, periode_txt, cabang_txt, total_unique, penyusun):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)

    m = ASSETS / "logo_madinah.png"
    f = ASSETS / "logo_mflash.png"
    if m.exists():
        s.shapes.add_picture(str(m), In(1.35), In(1.75), height=In(1.5))
    if f.exists():
        s.shapes.add_picture(str(f), In(1.45), In(3.75), height=In(1.9))

    _text(s, 4.7, 2.35, 8.2, 0.7, "LAPORAN KINERJA SERVICE",
          font=F_H, size=34, bold=True, color=NAVY)
    _text(s, 4.7, 3.05, 8.2, 0.55, periode_txt.upper(), font=F_H, size=22, color=NAVY_L)

    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(4.72), In(3.82), In(7.4), In(0.035))
    ln.fill.solid()
    ln.fill.fore_color.rgb = NAVY
    ln.line.fill.background()
    ln.shadow.inherit = False

    _text(s, 4.7, 4.12, 8.2, 0.4, cabang_txt.upper(), font=F_H, size=17, color=NAVY_L)
    _text(s, 4.7, 4.62, 8.2, 0.35,
          f"{nf(total_unique)} TRANSAKSI UNIK DIANALISIS", size=13, color=MUTED)
    if penyusun:
        _text(s, 4.7, 5.15, 8.2, 0.35, penyusun.upper(), font=F_H, size=15, bold=True, color=NAVY)
    _text(s, 4.7, 5.62, 8.2, 0.3,
          f"Dibuat otomatis dari dashboard · {date.today().strftime('%d %B %Y')}",
          size=10, color=MUTED)
    return s


def _slide_ringkasan(prs, d, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Ringkasan Kinerja")
    _text(s, 1.06, 0.94, 11.6, 0.3, meta["subjudul"], size=12, color=MUTED)

    tot = meta["total"]

    def pct(n):
        return f"{dec(n / tot * 100)}%" if tot else "0%"

    _kpi_row(s, [
        {"label": "TOTAL TRANSAKSI", "value": nf(tot), "sub": "sesuai filter aktif",
         "fill": NAVY, "text": WHITE, "subcolor": RGBColor(0xC7, 0xD3, 0xEA)},
        {"label": "SELESAI (DONE)", "value": nf(meta["done"]), "sub": f"{pct(meta['done'])} dari total",
         "fill": GREEN_B, "text": WHITE, "subcolor": RGBColor(0xDC, 0xF3, 0xE3)},
        {"label": "BATAL (CANCEL)", "value": nf(meta["cancel"]), "sub": f"{pct(meta['cancel'])} dari total",
         "fill": RED, "text": WHITE, "subcolor": RGBColor(0xF7, 0xDE, 0xDB)},
        {"label": "PENDING", "value": nf(meta["pending"]), "sub": f"{pct(meta['pending'])} dari total",
         "fill": AMBER, "text": WHITE, "subcolor": RGBColor(0xFB, 0xEC, 0xD5)},
        {"label": "RATA-RATA / HARI", "value": dec(meta["avg_day"]),
         "sub": f"{meta['period_days']} hari periode",
         "fill": BLUE, "text": WHITE, "subcolor": RGBColor(0xD9, 0xEE, 0xF9)},
    ], y=1.42, h=1.32)

    # grafik tren bulanan
    tren = meta["tren"]
    if tren is not None and len(tren):
        cd = CategoryChartData()
        cd.categories = [BULAN_NAMES[m][:3] for m in tren.index]
        cd.add_series("Transaksi", list(tren.values))
        gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(0.62), In(3.05), In(7.5), In(3.3), cd)
        _style_chart(gf.chart, colors=[NAVY], title="Volume transaksi per bulan",
                     val_max=None, cat_size=10)

    # highlight kanan
    _text(s, 8.35, 3.12, 4.37, 0.32, "Catatan Utama", font=F_H, size=15, bold=True, color=NAVY)
    yy = 3.58
    for t in meta["highlights"]:
        _bullet(s, 8.35, yy, 4.37, t, size=11.5, color=GREEN)
        yy += 0.78

    _footer(s, meta["footer"])
    return s


def _slide_status(prs, d, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Komposisi Status Pengerjaan")
    _text(s, 1.06, 0.94, 11.6, 0.3, meta["subjudul"], size=12, color=MUTED)

    tot = meta["total"]
    vals = [meta["done"], meta["cancel"], meta["pending"], meta["lainnya"]]
    labels = ["Done", "Cancel", "Pending", "Lainnya"]
    cols = [GREEN_B, RED, AMBER, RGBColor(0x94, 0xA3, 0xB8)]

    cd = CategoryChartData()
    cd.categories = [f"{l}  {dec(v / tot * 100) if tot else 0}%" for l, v in zip(labels, vals)]
    cd.add_series("Jumlah", vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, In(0.62), In(1.5), In(6.1), In(4.9), cd)
    ch = gf.chart
    _style_chart(ch, show_legend=True, legend_pos=XL_LEGEND_POSITION.BOTTOM,
                 show_labels=False, title="Proporsi hasil akhir transaksi")
    for i, pt in enumerate(ch.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = cols[i % len(cols)]

    # tabel ringkas kanan
    _text(s, 7.0, 1.52, 5.72, 0.34, "Rincian Angka", font=F_H, size=15, bold=True, color=NAVY)
    rows = [(labels[i], vals[i], cols[i]) for i in range(4)]
    for i, (nm, v, c) in enumerate(rows):
        y = 2.02 + i * 0.72
        _rect(s, 7.0, y, 5.72, 0.62, CARD if i % 2 == 0 else WHITE, line_color=LINE)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, In(7.2), In(y + 0.25), In(0.13), In(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()
        dot.shadow.inherit = False
        _text(s, 7.46, y, 2.3, 0.62, nm, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 9.7, y, 1.5, 0.62, nf(v), size=12, bold=True, color=NAVY,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 11.3, y, 1.2, 0.62, f"{dec(v / tot * 100) if tot else 0}%", size=12, bold=True,
              color=c, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    _rect(s, 7.0, 5.02, 5.72, 1.05, RGBColor(0xFF, 0xF8, 0xEC), line_color=RGBColor(0xF0, 0xD9, 0xA8))
    _text(s, 7.2, 5.12, 5.35, 0.9,
          "Pending adalah kondisi terkini, bukan akumulasi sepanjang periode. "
          "Transaksi lama yang sudah tuntas tidak lagi terhitung di sini.",
          size=10, color=RGBColor(0x7A, 0x5B, 0x18), anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, meta["footer"])
    return s


def _slide_cabang(prs, d, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Kinerja per Cabang")
    _text(s, 1.06, 0.94, 11.6, 0.3, "Diurutkan dari volume terbesar", size=12, color=MUTED)

    rekap = (d.groupby("CABANG")["STATUS_BUCKET"].value_counts().unstack(fill_value=0))
    for c in ["DONE", "PENDING", "CANCEL", "LAINNYA"]:
        if c not in rekap.columns:
            rekap[c] = 0
    rekap["TOTAL"] = rekap[["DONE", "PENDING", "CANCEL", "LAINNYA"]].sum(axis=1)
    rekap = rekap[rekap["TOTAL"] > 0].sort_values("TOTAL", ascending=False)
    rekap["DP"] = rekap["DONE"] / rekap["TOTAL"] * 100
    rekap["CP"] = rekap["CANCEL"] / rekap["TOTAL"] * 100

    items = list(rekap.itertuples())
    avg_done = (d["STATUS_BUCKET"] == "DONE").sum() / len(d) * 100 if len(d) else 0

    half = 8
    HDR = ["CABANG", "VOLUME", "% DONE", "% CANCEL"]
    CW = [2.0, 1.0, 1.0, 1.15]
    OFF = [0.0, 2.05, 3.15, 4.35]

    def draw(chunk, ox):
        xs = [ox + o for o in OFF]
        for i, h in enumerate(HDR):
            _text(s, xs[i], 1.44, CW[i], 0.3, h, size=8.5, bold=True, color=MUTED, spacing=0.8,
                  align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
        for i, r in enumerate(chunk):
            y = 1.8 + i * 0.53
            if i % 2 == 0:
                _rect(s, ox - 0.14, y - 0.03, 5.78, 0.5, CARD, shape=MSO_SHAPE.RECTANGLE)
            _text(s, xs[0], y, CW[0], 0.44, r.Index, size=10.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)
            _text(s, xs[1], y, CW[1], 0.44, nf(r.TOTAL), size=10.5, color=INK,
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            _text(s, xs[2], y, CW[2], 0.44, f"{dec(r.DP)}%", size=10.5, bold=True,
                  color=GREEN_B if r.DP >= avg_done + 2 else (RED if r.DP < avg_done - 3 else INK),
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            _text(s, xs[3], y, CW[3], 0.44, f"{dec(r.CP)}%", size=10.5, bold=True,
                  color=RED if r.CP >= 15 else (GREEN_B if r.CP <= 10 else INK),
                  align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    draw(items[:half], 0.74)
    if len(items) > half:
        draw(items[half:half * 2], 7.15)

    _text(s, 0.62, 6.28, 12.1, 0.3,
          f"Hijau = di atas rata-rata · Merah = perlu perhatian · Rata-rata penyelesaian: {dec(avg_done)}%",
          size=10, color=MUTED)
    if len(items) > half * 2:
        _text(s, 0.62, 6.58, 12.1, 0.3,
              f"Menampilkan {half*2} cabang dengan volume terbesar dari total {len(items)} cabang.",
              size=10, color=MUTED)
    _footer(s, meta["footer"])
    return s


def _slide_status_detail(prs, sub, meta, *, judul, warna, palette, jenis_func,
                          jenis_label, catatan):
    """Slide detail untuk satu status (Pending / Done / Cancel)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, judul)

    total_s = len(sub)
    tot_all = meta["total"]
    porsi = (total_s / tot_all * 100) if tot_all else 0
    _text(s, 1.06, 0.94, 11.6, 0.3,
          f"{nf(total_s)} transaksi — {dec(porsi)}% dari total pada filter aktif",
          size=12, color=MUTED)

    if total_s == 0:
        _text(s, 0.62, 3.0, 12.1, 0.5, "Tidak ada data untuk filter yang dipilih.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _footer(s, meta["footer"])
        return s

    sub = sub.copy()
    sub["JENIS"] = sub["STATUS PENGERJAAN"].apply(jenis_func)
    tek = sub[~sub["TEKNISI"].isin(["TIDAK ADA TEKNISI", "N/A"])]["TEKNISI"].value_counts()
    cab = sub["CABANG"].value_counts()
    ker = sub["KERUSAKAN"].value_counts()
    jns = sub["JENIS"].value_counts()
    pdays = _period_days(sub)

    _kpi_row(s, [
        {"label": "JUMLAH", "value": nf(total_s), "sub": f"{dec(porsi)}% dari total",
         "fill": warna, "text": WHITE, "subcolor": WHITE},
        {"label": "TEKNISI TERBANYAK",
         "value": (str(tek.index[0])[:22] if len(tek) else "-"),
         "sub": (f"{nf(tek.iloc[0])} unit ({dec(tek.iloc[0]/total_s*100)}%)" if len(tek) else "-")},
        {"label": "CABANG TERBANYAK", "value": (str(cab.index[0])[:18] if len(cab) else "-"),
         "sub": (f"{nf(cab.iloc[0])} unit ({dec(cab.iloc[0]/total_s*100)}%)" if len(cab) else "-")},
        {"label": "KERUSAKAN TERBANYAK", "value": (str(ker.index[0])[:20] if len(ker) else "-"),
         "sub": (f"{nf(ker.iloc[0])} unit ({dec(ker.iloc[0]/total_s*100)}%)" if len(ker) else "-")},
        {"label": "RATA-RATA / HARI", "value": dec(total_s / pdays, 2) if pdays else "0",
         "sub": f"{pdays} hari periode" if pdays else " "},
    ], y=1.42, h=1.3)

    # ranking teknisi (bar horizontal, terbesar di atas)
    top_t = tek.head(8)
    if len(top_t):
        cd = CategoryChartData()
        cd.categories = [str(i)[:20] for i in top_t.index[::-1]]
        cd.add_series("Jumlah", [int(v) for v in top_t.values[::-1]])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, In(0.62), In(2.96), In(6.0), In(3.35), cd)
        _style_chart(gf.chart, colors=[warna], title="Teknisi terbanyak",
                     val_max=int(top_t.iloc[0] * 1.25) + 1, cat_size=9)

    # ranking kerusakan
    top_k = ker.head(8)
    if len(top_k):
        cd = CategoryChartData()
        cd.categories = [str(i)[:20] for i in top_k.index[::-1]]
        cd.add_series("Jumlah", [int(v) for v in top_k.values[::-1]])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, In(6.85), In(2.96), In(5.87), In(3.35), cd)
        _style_chart(gf.chart, colors=[palette], title="Jenis kerusakan terbanyak",
                     val_max=int(top_k.iloc[0] * 1.25) + 1, cat_size=9)

    if catatan:
        _footer(s, catatan)
    else:
        _footer(s, meta["footer"])
    return s


def _slide_harian(prs, d, meta):
    """Rekap jumlah transaksi per hari + pola hari dalam pekan."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Rekap Transaksi Harian")
    _text(s, 1.06, 0.94, 11.6, 0.3, meta["subjudul"], size=12, color=MUTED)

    daily = meta["daily"]
    if daily is None or daily.empty:
        _text(s, 0.62, 3.0, 12.1, 0.5, "Tidak ada data tanggal untuk filter yang dipilih.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _footer(s, meta["footer"])
        return s

    hari_aktif = len(daily)
    rata = daily.mean()
    tmax, vmax = daily.idxmax(), int(daily.max())
    tmin, vmin = daily.idxmin(), int(daily.min())

    _kpi_row(s, [
        {"label": "HARI AKTIF", "value": nf(hari_aktif), "sub": "hari ada transaksi",
         "fill": NAVY, "text": WHITE, "subcolor": RGBColor(0xC7, 0xD3, 0xEA)},
        {"label": "RATA-RATA / HARI", "value": dec(rata), "sub": "transaksi per hari aktif",
         "fill": BLUE, "text": WHITE, "subcolor": RGBColor(0xD9, 0xEE, 0xF9)},
        {"label": "HARI TERTINGGI", "value": nf(vmax),
         "sub": f"{tgl_id(tmax)} · {HARI_NAMES[pd.Timestamp(tmax).dayofweek]}",
         "fill": GREEN_B, "text": WHITE, "subcolor": RGBColor(0xDC, 0xF3, 0xE3)},
        {"label": "HARI TERENDAH", "value": nf(vmin),
         "sub": f"{tgl_id(tmin)} · {HARI_NAMES[pd.Timestamp(tmin).dayofweek]}",
         "fill": AMBER, "text": WHITE, "subcolor": RGBColor(0xFB, 0xEC, 0xD5)},
        {"label": "SELISIH TERTINGGI", "value": f"+{dec((vmax/rata-1)*100, 0)}%",
         "sub": "di atas rata-rata harian",
         "fill": RED, "text": WHITE, "subcolor": RGBColor(0xF7, 0xDE, 0xDB)},
    ], y=1.42, h=1.3)

    # --- grafik harian ---
    # kalau periodenya panjang, label sumbu-X diringkas agar tidak menumpuk
    n = len(daily)
    step = max(1, n // 14)
    cats = []
    for i, dt in enumerate(daily.index):
        cats.append(tgl_id(dt, pendek=True) if i % step == 0 else " " * (i + 1))

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Transaksi", [int(v) for v in daily.values])
    ctype = XL_CHART_TYPE.LINE if n > 40 else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = s.shapes.add_chart(ctype, In(0.62), In(3.0), In(8.4), In(3.35), cd)
    ch = gf.chart
    _style_chart(ch, colors=None if ctype == XL_CHART_TYPE.LINE else [NAVY],
                 title="Jumlah transaksi per hari",
                 show_labels=(n <= 20), cat_size=8,
                 val_max=int(vmax * 1.15) + 1)
    if ctype == XL_CHART_TYPE.LINE:
        ser = ch.series[0]
        ser.format.line.color.rgb = NAVY
        ser.format.line.width = Pt(1.75)
        ser.smooth = False

    # --- pola hari dalam pekan ---
    dow = meta["dow_avg"]
    _text(s, 9.25, 3.05, 3.47, 0.32, "Rata-rata per Hari", font=F_H, size=14, bold=True, color=NAVY)
    if dow is not None and len(dow):
        mx = dow.max()
        for i, (idx, val) in enumerate(dow.items()):
            y = 3.5 + i * 0.42
            is_top = val == mx
            col = GREEN_B if is_top else NAVY_L
            _text(s, 9.25, y, 0.95, 0.3, HARI_NAMES[idx], size=10.5,
                  bold=is_top, color=INK, anchor=MSO_ANCHOR.MIDDLE)
            _rect(s, 10.25, y + 0.08, 1.75, 0.16, RGBColor(0xE6, 0xEA, 0xF2), radius=0.4)
            _rect(s, 10.25, y + 0.08, max(0.05, 1.75 * (val / mx)), 0.16, col, radius=0.4)
            _text(s, 12.05, y, 0.67, 0.3, dec(val, 0), size=10.5, bold=is_top,
                  color=col, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, meta["footer"])
    return s


def _slide_hari_tertinggi(prs, d, meta):
    """Breakdown tanggal dengan transaksi tertinggi."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Hari dengan Transaksi Tertinggi")

    daily = meta["daily"]
    if daily is None or daily.empty:
        _text(s, 0.62, 3.0, 12.1, 0.5, "Tidak ada data tanggal untuk filter yang dipilih.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _footer(s, meta["footer"])
        return s

    rata = daily.mean()
    top = daily.sort_values(ascending=False).head(10)
    tmax, vmax = top.index[0], int(top.iloc[0])

    _text(s, 1.06, 0.94, 11.6, 0.3,
          f"Puncak tertinggi: {tgl_id(tmax)} ({HARI_NAMES[pd.Timestamp(tmax).dayofweek]}) "
          f"dengan {nf(vmax)} transaksi — {dec((vmax/rata-1)*100, 0)}% di atas rata-rata harian "
          f"({dec(rata, 0)} transaksi).",
          size=12, color=MUTED)

    # --- tabel 10 tanggal tertinggi ---
    _text(s, 0.62, 1.5, 6.5, 0.32, "10 Tanggal Tertinggi", font=F_H, size=14, bold=True, color=NAVY)
    HDR = [("#", 0.0, 0.4, PP_ALIGN.LEFT), ("TANGGAL", 0.42, 1.85, PP_ALIGN.LEFT),
           ("HARI", 2.3, 1.15, PP_ALIGN.LEFT), ("JUMLAH", 3.5, 1.05, PP_ALIGN.RIGHT),
           ("VS RATA-RATA", 4.6, 1.55, PP_ALIGN.RIGHT)]
    for h, ox, w, al in HDR:
        _text(s, 0.74 + ox, 1.9, w, 0.28, h, size=8.5, bold=True, color=MUTED,
              spacing=0.8, align=al)

    for i, (dt, v) in enumerate(top.items()):
        y = 2.24 + i * 0.42
        if i % 2 == 0:
            _rect(s, 0.62, y - 0.02, 6.0, 0.4, CARD, shape=MSO_SHAPE.RECTANGLE)
        v = int(v)
        delta = (v / rata - 1) * 100
        dow_i = pd.Timestamp(dt).dayofweek
        akhir_pekan = dow_i >= 5
        _text(s, 0.74, y, 0.4, 0.36, str(i + 1), size=10,
              bold=(i == 0), color=NAVY if i == 0 else MUTED, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 1.16, y, 1.85, 0.36, tgl_id(dt), size=10, bold=(i == 0),
              color=INK, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 3.04, y, 1.15, 0.36, HARI_NAMES[dow_i], size=10,
              color=GREEN if akhir_pekan else MUTED,
              bold=akhir_pekan, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 4.24, y, 1.05, 0.36, nf(v), size=10, bold=True, color=NAVY,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 5.34, y, 1.55, 0.36, f"+{dec(delta, 0)}%", size=10, bold=True,
              color=GREEN_B, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # --- grafik 10 tanggal tertinggi ---
    cd = CategoryChartData()
    rev = list(top.items())[::-1]
    cd.categories = [f"{tgl_id(dt, pendek=True)} ({HARI_NAMES[pd.Timestamp(dt).dayofweek][:3]})"
                     for dt, _ in rev]
    cd.add_series("Transaksi", [int(v) for _, v in rev])
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, In(7.0), In(1.5), In(5.72), In(4.0), cd)
    _style_chart(gf.chart, colors=[NAVY], title="Perbandingan 10 tanggal tertinggi",
                 val_max=int(vmax * 1.2) + 1, cat_size=9)

    # --- catatan pola ---
    dow_avg = meta["dow_avg"]
    catatan = ""
    if dow_avg is not None and len(dow_avg):
        top_dow = dow_avg.idxmax()
        low_dow = dow_avg.idxmin()
        selisih = (dow_avg.max() / dow_avg.min() - 1) * 100
        akhir_pekan_top = sum(1 for dt in top.index if pd.Timestamp(dt).dayofweek >= 5)
        catatan = (f"{HARI_NAMES[top_dow]} rata-rata tersibuk ({dec(dow_avg.max(), 0)} transaksi/hari), "
                   f"{dec(selisih, 0)}% di atas {HARI_NAMES[low_dow]} yang paling sepi "
                   f"({dec(dow_avg.min(), 0)}). Dari 10 tanggal tertinggi, {akhir_pekan_top} di antaranya "
                   f"jatuh pada akhir pekan.")

    if catatan:
        _rect(s, 7.0, 5.66, 5.72, 0.98, RGBColor(0xF2, 0xF8, 0xF3),
              line_color=RGBColor(0xC6, 0xE2, 0xCB))
        _text(s, 7.2, 5.74, 5.35, 0.84, catatan, size=10,
              color=RGBColor(0x2C, 0x5F, 0x2D), anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, meta["footer"])
    return s


def _slide_banding_bulan(prs, d, sf, meta):
    """Bandingkan bulan berjalan vs bulan sebelumnya secara setara.

    Bulan berjalan biasanya belum penuh, jadi bulan sebelumnya ikut dipotong
    pada tanggal yang sama (proporsi hari yang sama) agar adil. Rata-rata
    harian tetap ditampilkan sebagai pembanding kedua.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Bulan Berjalan vs Bulan Sebelumnya")

    tgl = d["TGL PENGIRIMAN"].dropna() if d is not None and not d.empty else pd.Series(dtype='datetime64[ns]')
    if tgl.empty:
        _text(s, 0.62, 3.0, 12.1, 0.5, "Data tanggal tidak tersedia.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _footer(s, meta["footer"])
        return s

    tmax = tgl.max()
    th, bl = tmax.year, tmax.month

    # Hari terakhir sering belum lengkap (data ditarik di tengah hari). Kalau
    # volumenya jauh di bawah kebiasaan, hari itu dibuang supaya tidak menyeret
    # turun angka bulan berjalan.
    _hb = tgl[(tgl.dt.year == th) & (tgl.dt.month == bl)].dt.day.value_counts().sort_index()
    hari_n = int(tmax.day)
    hari_dibuang = None
    if len(_hb) >= 3 and _hb.get(hari_n, 0) < 0.4 * _hb.iloc[:-1].median():
        hari_dibuang = hari_n
        hari_n -= 1

    bl_prev, th_prev = (bl - 1, th) if bl > 1 else (12, th - 1)

    def potong(df, kol, bulan, tahun):
        t = df[kol]
        return df[(t.dt.year == tahun) & (t.dt.month == bulan) & (t.dt.day <= hari_n)]

    cur = potong(d, "TGL PENGIRIMAN", bl, th)
    prev = potong(d, "TGL PENGIRIMAN", bl_prev, th_prev)

    def stat(x):
        if x.empty:
            return dict(total=0, done=0, cancel=0, pending=0)
        vc = x["STATUS_BUCKET"].value_counts()
        return dict(total=len(x), done=int(vc.get("DONE", 0)),
                    cancel=int(vc.get("CANCEL", 0)), pending=int(vc.get("PENDING", 0)))

    a, b = stat(prev), stat(cur)

    nama_ini = f"{BULAN_NAMES[bl]} {th}"
    nama_lalu = f"{BULAN_NAMES[bl_prev]} {th_prev}"
    _text(s, 1.06, 0.94, 11.6, 0.32,
          f"Dibandingkan setara: tanggal 1–{hari_n} pada kedua bulan "
          f"({nama_lalu} vs {nama_ini})", size=12, color=MUTED)

    def delta(baru, lama):
        if not lama:
            return "—", MUTED
        p = (baru - lama) / lama * 100
        warna = GREEN_B if p > 0 else (RED if p < 0 else MUTED)
        return f"{'+' if p > 0 else ''}{dec(p)}%", warna

    dt_total, w_total = delta(b["total"], a["total"])
    dt_done, w_done = delta(b["done"], a["done"])
    dt_cancel, w_cancel = delta(b["cancel"], a["cancel"])

    kartu = [
        {"label": f"TRANSAKSI 1–{hari_n}", "value": nf(b["total"]),
         "sub": f"{nama_lalu}: {nf(a['total'])}  ({dt_total})",
         "fill": NAVY, "text": WHITE, "subcolor": RGBColor(0xC7, 0xD3, 0xEA)},
        {"label": "SELESAI (DONE)", "value": nf(b["done"]),
         "sub": f"{nama_lalu}: {nf(a['done'])}  ({dt_done})",
         "fill": GREEN_B, "text": WHITE, "subcolor": RGBColor(0xDC, 0xF3, 0xE3)},
        {"label": "BATAL (CANCEL)", "value": nf(b["cancel"]),
         "sub": f"{nama_lalu}: {nf(a['cancel'])}  ({dt_cancel})",
         "fill": RED, "text": WHITE, "subcolor": RGBColor(0xF7, 0xDE, 0xDB)},
        {"label": "RATA-RATA / HARI", "value": dec(b["total"] / hari_n if hari_n else 0),
         "sub": f"{nama_lalu}: {dec(a['total'] / hari_n if hari_n else 0)}",
         "fill": BLUE, "text": WHITE, "subcolor": RGBColor(0xD9, 0xEE, 0xF9)},
    ]

    # tambah omzet & laba bila data penjualan tersedia
    om_a = om_b = lb_a = lb_b = None
    if sf is not None and not sf.empty and "TGL" in sf.columns:
        scur = potong(sf, "TGL", bl, th)
        sprev = potong(sf, "TGL", bl_prev, th_prev)
        om_a, om_b = sprev["TOTAL HARGA"].sum(), scur["TOTAL HARGA"].sum()
        lb_a, lb_b = sprev["LABA"].sum(), scur["LABA"].sum()
        dt_om, _ = delta(om_b, om_a)
        dt_lb, _ = delta(lb_b, lb_a)
        kartu.append({"label": "OMZET", "value": rp(om_b),
                      "sub": f"{nama_lalu}: {rp(om_a)}  ({dt_om})",
                      "fill": RGBColor(0x7C, 0x3A, 0xED), "text": WHITE,
                      "subcolor": RGBColor(0xE4, 0xD9, 0xFA)})
        kartu.append({"label": "LABA KOTOR", "value": rp(lb_b),
                      "sub": f"{nama_lalu}: {rp(lb_a)}  ({dt_lb})",
                      "fill": ORANGE, "text": WHITE,
                      "subcolor": RGBColor(0xFB, 0xEC, 0xD5)})

    _kpi_row(s, kartu, y=1.5, h=1.35)

    # Grafik perbandingan status — Pending sengaja TIDAK disertakan karena
    # angkanya kondisi terkini, bukan kejadian pada bulan itu: pendingan bulan
    # lalu sebagian sudah tuntas sehingga tampak selalu jauh lebih kecil.
    cd = CategoryChartData()
    cd.categories = ["Total", "Done", "Cancel"]
    cd.add_series(nama_lalu, [a["total"], a["done"], a["cancel"]])
    cd.add_series(nama_ini, [b["total"], b["done"], b["cancel"]])
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(0.62), In(3.15),
                            In(6.6), In(3.2), cd)
    ch = gf.chart
    _style_chart(ch, title=f"Perbandingan status (tanggal 1–{hari_n})",
                 show_legend=True, legend_pos=XL_LEGEND_POSITION.TOP,
                 colors=[RGBColor(0xA9, 0xB6, 0xD8), NAVY], cat_size=10,
                 val_max=int(max(a["total"], b["total"]) * 1.25) + 1)

    # grafik harian kumulatif
    def harian(x, kol):
        if x.empty:
            return [0] * hari_n
        g = x.groupby(x[kol].dt.day).size()
        return [int(g.get(i, 0)) for i in range(1, hari_n + 1)]

    cd2 = CategoryChartData()
    cd2.categories = [str(i) for i in range(1, hari_n + 1)]
    cd2.add_series(nama_lalu, harian(prev, "TGL PENGIRIMAN"))
    cd2.add_series(nama_ini, harian(cur, "TGL PENGIRIMAN"))
    gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE, In(7.45), In(3.15), In(5.27), In(3.2), cd2)
    ch2 = gf2.chart
    _style_chart(ch2, title="Transaksi harian (tanggal 1–%d)" % hari_n,
                 show_legend=True, legend_pos=XL_LEGEND_POSITION.TOP,
                 show_labels=False, cat_size=7)
    for i, col in enumerate([RGBColor(0xA9, 0xB6, 0xD8), NAVY]):
        ch2.series[i].format.line.color.rgb = col
        ch2.series[i].format.line.width = Pt(2)
        ch2.series[i].smooth = False

    # catatan bawah
    hari_penuh = calendar.monthrange(th_prev, bl_prev)[1]
    proyeksi = ""
    if hari_n < calendar.monthrange(th, bl)[1] and b["total"]:
        proy = b["total"] / hari_n * calendar.monthrange(th, bl)[1]
        proyeksi = (f" Bila laju bertahan, {nama_ini} diperkirakan menutup di sekitar "
                    f"{nf(round(proy))} transaksi.")
    buang = ""
    if hari_dibuang:
        buang = (f" Tanggal {hari_dibuang} dikecualikan karena datanya belum lengkap "
                 f"sehari penuh.")
    _text(s, 0.62, 6.62, 12.1, 0.34,
          "Pending tidak ikut dibandingkan: angkanya kondisi terkini, bukan kejadian "
          "bulan tersebut — pendingan bulan lalu sebagian sudah tuntas sehingga selalu "
          "tampak jauh lebih kecil.", size=9.5, color=RGBColor(0x7A, 0x5B, 0x18))
    _footer(s, f"{nama_ini} baru berjalan {hari_n} hari, sehingga {nama_lalu} "
               f"(sebenarnya {hari_penuh} hari) ikut dipotong sampai tanggal {hari_n} "
               f"agar perbandingannya setara.{buang}{proyeksi}")
    return s


def _slide_penjualan(prs, sf, meta):
    """Rekap penjualan: omzet, modal (harga beli), dan laba kotor."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Penjualan — Modal, Omzet & Laba")

    if sf is None or sf.empty:
        _text(s, 1.06, 0.94, 11.6, 0.3, meta["subjudul"], size=12, color=MUTED)
        _text(s, 0.62, 3.0, 12.1, 0.5,
              "Tidak ada data penjualan pada periode/cabang ini.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _footer(s, meta["footer"])
        return s

    omzet = sf['TOTAL HARGA'].sum()
    modal = sf['MODAL'].sum()
    laba = sf['LABA'].sum()
    margin = (laba / omzet * 100) if omzet else 0
    n_faktur = sf['NO FAKTUR'].nunique()
    qty = sf['QTY'].sum()

    _text(s, 1.06, 0.94, 11.6, 0.3,
          f"{meta['subjudul']} · {nf(n_faktur)} faktur · {nf(qty)} unit terjual",
          size=12, color=MUTED)

    _kpi_row(s, [
        {"label": "OMZET (HARGA JUAL)", "value": rp(omzet), "sub": f"{nf(len(sf))} baris",
         "fill": NAVY, "text": WHITE, "subcolor": RGBColor(0xC7, 0xD3, 0xEA)},
        {"label": "MODAL (HARGA BELI)", "value": rp(modal),
         "sub": f"{dec(modal/omzet*100 if omzet else 0)}% dari omzet",
         "fill": RED, "text": WHITE, "subcolor": RGBColor(0xF7, 0xDE, 0xDB)},
        {"label": "LABA KOTOR", "value": rp(laba), "sub": f"margin {dec(margin)}%",
         "fill": GREEN_B, "text": WHITE, "subcolor": RGBColor(0xDC, 0xF3, 0xE3)},
        {"label": "RATA-RATA / FAKTUR", "value": rp(omzet / n_faktur if n_faktur else 0),
         "sub": f"laba {rp(laba/n_faktur if n_faktur else 0)}/faktur",
         "fill": BLUE, "text": WHITE, "subcolor": RGBColor(0xD9, 0xEE, 0xF9)},
        {"label": "LABA / UNIT", "value": rp(laba / qty if qty else 0),
         "sub": f"dari {nf(qty)} unit",
         "fill": ORANGE, "text": WHITE, "subcolor": RGBColor(0xFB, 0xEC, 0xD5)},
    ], y=1.42, h=1.3)

    kat = sf['KATEGORI BARANG'].astype(str).str.strip().str.upper()
    gk = (sf.assign(KAT=kat).groupby('KAT')
          .agg(omzet=('TOTAL HARGA', 'sum'), modal=('MODAL', 'sum'),
               laba=('LABA', 'sum'))
          .sort_values('omzet', ascending=False))

    # grafik: modal vs laba per kategori (6 teratas)
    top = gk.head(6)[::-1]
    if len(top):
        cd = CategoryChartData()
        cd.categories = [str(i)[:16] for i in top.index]
        cd.add_series("Modal", [float(v) for v in top['modal']])
        cd.add_series("Laba", [float(v) for v in top['laba']])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, In(0.62), In(2.96),
                                In(6.4), In(3.35), cd)
        ch = gf.chart
        _style_chart(ch, title="Modal & laba per kategori barang (panjang = omzet)",
                     show_legend=True, legend_pos=XL_LEGEND_POSITION.TOP,
                     show_labels=False, cat_size=9)
        for i, col in enumerate([RED, GREEN_B]):
            ch.series[i].format.fill.solid()
            ch.series[i].format.fill.fore_color.rgb = col
        # angka sumbu bernilai miliaran bikin label menumpuk; nilai pastinya
        # sudah tersedia di tabel sebelah kanan
        try:
            from pptx.enum.chart import XL_TICK_LABEL_POSITION
            ch.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
        except Exception:
            pass

    # tabel margin per kategori
    _text(s, 7.25, 3.0, 5.47, 0.32, "Margin per kategori", font=F_H, size=14,
          bold=True, color=NAVY)
    yy = 3.44
    for j, (h, x, w, al) in enumerate([("KATEGORI", 7.25, 1.7, PP_ALIGN.LEFT),
                                        ("OMZET", 8.95, 1.25, PP_ALIGN.RIGHT),
                                        ("LABA", 10.25, 1.25, PP_ALIGN.RIGHT),
                                        ("MARGIN", 11.55, 1.17, PP_ALIGN.RIGHT)]):
        _text(s, x, yy, w, 0.26, h, size=8.5, bold=True, color=MUTED,
              spacing=0.8, align=al)
    yy += 0.3
    for i, (lbl, row) in enumerate(gk.head(6).iterrows()):
        if i % 2 == 0:
            _rect(s, 7.12, yy - 0.02, 5.6, 0.42, CARD, shape=MSO_SHAPE.RECTANGLE)
        mg = (row['laba'] / row['omzet'] * 100) if row['omzet'] else 0
        _text(s, 7.25, yy, 1.7, 0.38, str(lbl)[:16], size=10, color=INK,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 8.95, yy, 1.25, 0.38, rp(row['omzet']), size=9.5, color=INK,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 10.25, yy, 1.25, 0.38, rp(row['laba']), size=9.5, color=INK,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 11.55, yy, 1.17, 0.38, f"{dec(mg)}%", size=10, bold=True,
              color=(GREEN_B if mg >= 30 else RED if mg < 12 else INK),
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.44

    # catatan: jasa bermodal nol
    jasa_omz = gk.loc['JASA', 'omzet'] if 'JASA' in gk.index else 0
    porsi_jasa = (jasa_omz / omzet * 100) if omzet else 0
    _footer(s, f"Modal diambil dari kolom HARGA BELI yang sudah berupa total per baris. "
               f"Kategori JASA hampir seluruhnya bermodal nol sehingga marginnya tampil "
               f"100% ({dec(porsi_jasa)}% dari omzet) — biaya tenaga kerja tidak "
               f"dibebankan per faktur. Laba di sini laba kotor, belum dikurangi biaya operasional.")
    return s


def _slide_mlf(prs, mlf, meta):
    """Rekap penjualan Voucher Tiket MLF semua cabang."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Voucher Tiket MLF")

    if mlf is None or mlf.empty:
        _text(s, 1.06, 0.94, 11.6, 0.3, meta["subjudul"], size=12, color=MUTED)
        _text(s, 0.62, 3.0, 12.1, 0.5,
              "Tidak ada penjualan voucher MLF pada periode/cabang ini.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _footer(s, meta["footer"])
        return s

    qty = mlf['QTY'].sum()
    omzet = mlf['TOTAL HARGA'].sum()
    modal = mlf['MODAL'].sum()
    laba = mlf['LABA'].sum()
    margin = (laba / omzet * 100) if omzet else 0
    n_cab = mlf['CABANG'].nunique()
    hari = mlf['TGL'].dt.normalize().nunique()

    _text(s, 1.06, 0.94, 11.6, 0.3,
          f"{meta['subjudul']} · {nf(qty)} voucher terjual di {n_cab} cabang",
          size=12, color=MUTED)

    _kpi_row(s, [
        {"label": "VOUCHER TERJUAL", "value": nf(qty), "sub": f"{nf(len(mlf))} transaksi",
         "fill": NAVY, "text": WHITE, "subcolor": RGBColor(0xC7, 0xD3, 0xEA)},
        {"label": "OMZET", "value": rp(omzet),
         "sub": f"rata-rata {rp(omzet/qty if qty else 0)}/voucher",
         "fill": BLUE, "text": WHITE, "subcolor": RGBColor(0xD9, 0xEE, 0xF9)},
        {"label": "MODAL", "value": rp(modal),
         "sub": f"{dec(modal/omzet*100 if omzet else 0)}% dari omzet",
         "fill": RED, "text": WHITE, "subcolor": RGBColor(0xF7, 0xDE, 0xDB)},
        {"label": "LABA KOTOR", "value": rp(laba), "sub": f"margin {dec(margin)}%",
         "fill": GREEN_B, "text": WHITE, "subcolor": RGBColor(0xDC, 0xF3, 0xE3)},
        {"label": "RATA-RATA / HARI", "value": dec(qty / hari if hari else 0),
         "sub": f"{hari} hari ada penjualan",
         "fill": ORANGE, "text": WHITE, "subcolor": RGBColor(0xFB, 0xEC, 0xD5)},
    ], y=1.42, h=1.3)

    # ranking cabang
    pc = mlf.groupby('CABANG')['QTY'].sum().sort_values(ascending=False).head(10)
    if len(pc):
        rev = pc[::-1]
        cd = CategoryChartData()
        cd.categories = list(rev.index)
        cd.add_series("Voucher", [int(v) for v in rev.values])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, In(0.62), In(2.96),
                                In(6.0), In(3.35), cd)
        _style_chart(gf.chart, colors=[NAVY], title="Cabang dengan penjualan terbanyak",
                     val_max=int(pc.iloc[0] * 1.25) + 1, cat_size=9)

    # tren per bulan
    bl = mlf.groupby(mlf['TGL'].dt.month)['QTY'].sum().sort_index()
    if len(bl):
        cd = CategoryChartData()
        cd.categories = [BULAN_NAMES[int(b)][:3] for b in bl.index]
        cd.add_series("Voucher", [int(v) for v in bl.values])
        gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(6.85), In(2.96),
                                In(5.87), In(3.35), cd)
        _style_chart(gf.chart, colors=[BLUE], title="Penjualan voucher per bulan",
                     val_max=int(bl.max() * 1.25) + 1, cat_size=10)

    top_cab = pc.index[0] if len(pc) else "-"
    top_val = int(pc.iloc[0]) if len(pc) else 0
    _footer(s, f"Cabang terbanyak: {top_cab} ({nf(top_val)} voucher, "
               f"{dec(top_val/qty*100 if qty else 0)}% dari total). "
               f"Periode data: {mlf['TGL'].min():%d/%m/%Y} – {mlf['TGL'].max():%d/%m/%Y}.")
    return s


def _slide_bagihasil(prs, jasa, meta, *, tarif_map, tarif_flat, prioritas, periode_txt):
    """Omzet jasa & bagi hasil per teknisi."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Bagi Hasil Teknisi")

    if jasa is None or jasa.empty:
        _text(s, 1.06, 0.94, 11.6, 0.3, meta["subjudul"], size=12, color=MUTED)
        _text(s, 0.62, 3.0, 12.1, 0.5,
              "Tidak ada transaksi jasa pada periode/cabang ini.",
              size=16, color=MUTED, align=PP_ALIGN.CENTER)
        _footer(s, meta["footer"])
        return s

    omzet = jasa['TOTAL HARGA'].sum()
    bh = jasa['BAGI_HASIL'].sum()
    fl = jasa['FLAT'].sum()
    selisih = bh - fl
    n_tek = jasa.loc[jasa['TEKNISI'] != 'TIDAK ADA TEKNISI', 'TEKNISI'].nunique()

    _text(s, 1.06, 0.94, 11.6, 0.3,
          f"Periode penggajian: {periode_txt} · {n_tek} teknisi",
          size=12, color=MUTED)

    _kpi_row(s, [
        {"label": "OMZET JASA", "value": rp(omzet), "sub": f"{nf(len(jasa))} baris",
         "fill": NAVY, "text": WHITE, "subcolor": RGBColor(0xC7, 0xD3, 0xEA)},
        {"label": "BAGI HASIL (ATURAN)", "value": rp(bh),
         "sub": f"{dec(bh/omzet*100 if omzet else 0)}% dari omzet jasa",
         "fill": GREEN_B, "text": WHITE, "subcolor": RGBColor(0xDC, 0xF3, 0xE3)},
        {"label": f"PEMBANDING FLAT {tarif_flat:.0f}%", "value": rp(fl),
         "sub": f"omzet jasa × {tarif_flat:.0f}%",
         "fill": RGBColor(0x7C, 0x3A, 0xED), "text": WHITE,
         "subcolor": RGBColor(0xE4, 0xD9, 0xFA)},
        {"label": "SELISIH", "value": rp(selisih),
         "sub": ("aturan lebih besar" if selisih > 0
                 else "flat lebih besar" if selisih < 0 else "sama"),
         "fill": (ORANGE if selisih >= 0 else RED), "text": WHITE,
         "subcolor": RGBColor(0xFB, 0xEC, 0xD5)},
        {"label": "RATA-RATA / TEKNISI", "value": rp(bh / n_tek if n_tek else 0),
         "sub": f"dari {n_tek} teknisi",
         "fill": RGBColor(0x0F, 0x8A, 0x82), "text": WHITE,
         "subcolor": RGBColor(0xD5, 0xEE, 0xEC)},
    ], y=1.42, h=1.3)

    # 10 teknisi teratas (nama + cabang)
    rek = (jasa[jasa['TEKNISI'] != 'TIDAK ADA TEKNISI']
           .groupby(['TEKNISI', 'CABANG'])['BAGI_HASIL'].sum()
           .sort_values(ascending=False).head(10))
    if len(rek):
        rev = rek[::-1]
        cats = [f"{t[:18]} — {c[:9]}" for t, c in rev.index]
        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series("Bagi Hasil", [float(v) for v in rev.values])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, In(0.62), In(2.96),
                                In(6.9), In(3.35), cd)
        _style_chart(gf.chart, colors=[GREEN_B],
                     title="10 teknisi dengan bagi hasil tertinggi",
                     val_max=float(rek.iloc[0] * 1.3), cat_size=8,
                     number_format='#,##0')

    # komposisi per kategori tarif
    gt = jasa.groupby('TARIF_LABEL').agg(
        omzet=('TOTAL HARGA', 'sum'), bh=('BAGI_HASIL', 'sum')).sort_values(
        'omzet', ascending=False)
    _text(s, 7.75, 3.0, 4.97, 0.32, "Komposisi menurut tarif", font=F_H, size=14,
          bold=True, color=NAVY)
    yy = 3.44
    _text(s, 7.75, yy, 1.75, 0.26, "KATEGORI", size=8.5, bold=True, color=MUTED, spacing=0.8)
    _text(s, 9.5, yy, 0.75, 0.26, "TARIF", size=8.5, bold=True, color=MUTED,
          spacing=0.8, align=PP_ALIGN.RIGHT)
    _text(s, 10.3, yy, 1.2, 0.26, "OMZET", size=8.5, bold=True, color=MUTED,
          spacing=0.8, align=PP_ALIGN.RIGHT)
    _text(s, 11.55, yy, 1.17, 0.26, "BAGI HASIL", size=8.5, bold=True, color=MUTED,
          spacing=0.8, align=PP_ALIGN.RIGHT)
    yy += 0.3
    for i, (lbl, row) in enumerate(gt.iterrows()):
        if i >= 6:
            break
        if i % 2 == 0:
            _rect(s, 7.62, yy - 0.02, 5.1, 0.42, CARD, shape=MSO_SHAPE.RECTANGLE)
        tr = tarif_map.get(lbl, 0.0)
        _text(s, 7.75, yy, 1.75, 0.38, str(lbl), size=10, color=INK,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 9.5, yy, 0.75, 0.38, f"{tr:.0f}%", size=10, bold=True, color=NAVY,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 10.3, yy, 1.2, 0.38, rp(row['omzet']), size=9.5, color=INK,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 11.55, yy, 1.17, 0.38, rp(row['bh']), size=9.5, bold=True, color=GREEN,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.44

    ket = " · ".join(f"{k} {v:.0f}%" for k, v in tarif_map.items())
    _footer(s, f"Tarif: {ket}. Prioritas bila dua kata kunci: {prioritas}. "
               f"Cutoff penggajian tanggal 24 s/d 23. Angka berbasis omzet jasa, "
               f"belum dikurangi biaya apa pun.")
    return s


def _slide_penutup(prs, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_image(s, prs)
    _logos(s)
    _title(s, "Kesimpulan & Tindak Lanjut")
    _text(s, 1.06, 0.94, 11.6, 0.3, meta["subjudul"], size=12, color=MUTED)

    for i, (h, dsc) in enumerate(meta["kesimpulan"][:4]):
        col, row = i % 2, i // 2
        x, y = 0.62 + col * 6.22, 1.62 + row * 2.42
        _rect(s, x, y, 5.86, 2.06, WHITE, line_color=LINE)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x + 0.28), In(y + 0.3), In(0.44), In(0.44))
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        badge.shadow.inherit = False
        _text(s, x + 0.28, y + 0.3, 0.44, 0.44, str(i + 1), font=F_H, size=15, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.9, y + 0.3, 4.7, 0.46, h, size=13, bold=True, color=NAVY,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.3, y + 0.92, 5.28, 1.1, dsc, size=10.5, color=INK)

    _footer(s, meta["footer"])
    return s


# =====================================================================
# ENTRY POINT
# =====================================================================
def build_deck(df_filtered, *, total_unique_all, total_raw_rows,
               f_tahun="Semua Tahun", f_bulan="Semua Bulan", f_cabang="Semua Cabang",
               penyusun="", sertakan=("ringkasan", "status", "banding_bulan", "harian",
                                       "hari_tertinggi", "cabang", "pending", "done",
                                       "cancel", "penjualan", "mlf", "bagihasil",
                                       "penutup"),
               sales_filtered=None, tarif_map=None, tarif_flat=30.0,
               prioritas="Normal", periode_bagihasil=None, tahun_fokus="auto"):
    """Bangun deck presentasi.

    Parameter tambahan:
      sales_filtered  : data faktur penjualan yang sudah difilter (untuk slide
                        Voucher MLF dan Bagi Hasil Teknisi). Boleh None.
      tarif_map       : dict tarif bagi hasil dalam persen, mis.
                        {'Interface':20,'Normal':30,'Mati Total':32,'Promo':60,
                         'Lainnya':30}
      tarif_flat      : tarif pembanding (persen)
      prioritas       : kata kunci yang menang bila satu nama mengandung dua
      periode_bagihasil : (tahun, bulan_gaji) untuk cutoff 24->23; None = ikut filter
      tahun_fokus     : batasi seluruh isi deck ke satu tahun. "auto" (bawaan)
                        memakai tahun terbaru yang ada di data; None = tanpa batas.
    """
    d = df_filtered

    # --- batasi ke satu tahun saja (bawaan: tahun terbaru, mis. 2026) ---
    if tahun_fokus is not None and not d.empty and d["TAHUN"].notna().any():
        if tahun_fokus == "auto":
            tahun_fokus = int(d["TAHUN"].dropna().max())
        else:
            tahun_fokus = int(tahun_fokus)
        d = d[d["TAHUN"] == tahun_fokus]
        if f_tahun == "Semua Tahun":
            f_tahun = tahun_fokus
        if sales_filtered is not None and not sales_filtered.empty:
            _st = sales_filtered
            if "TAHUN" in _st.columns:
                sales_filtered = _st[_st["TAHUN"] == tahun_fokus]
            elif "TGL FAKTUR" in _st.columns:
                _tg = pd.to_datetime(_st["TGL FAKTUR"], errors="coerce")
                sales_filtered = _st[_tg.dt.year == tahun_fokus]

    prs = Presentation()
    prs.slide_width = In(13.333)
    prs.slide_height = In(7.5)

    total = len(d)
    vc = d["STATUS_BUCKET"].value_counts()
    done = int(vc.get("DONE", 0))
    cancel = int(vc.get("CANCEL", 0))
    pending = int(vc.get("PENDING", 0))
    lainnya = int(vc.get("LAINNYA", 0))
    pdays = _period_days(d)
    avg_day = (total / pdays) if pdays else 0

    # teks periode
    if f_tahun == "Semua Tahun" and f_bulan == "Semua Bulan":
        periode = "Seluruh Periode Data"
    elif f_bulan == "Semua Bulan":
        periode = f"Tahun {f_tahun}"
    elif f_tahun == "Semua Tahun":
        periode = f"Bulan {BULAN_NAMES[int(f_bulan)]} (semua tahun)"
    else:
        periode = f"{BULAN_NAMES[int(f_bulan)]} {f_tahun}"
    cabang_txt = "Seluruh Cabang" if f_cabang == "Semua Cabang" else f"Cabang {f_cabang}"
    subjudul = f"{periode} · {cabang_txt}"

    # tren bulanan
    tren = None
    if total and d["BULAN"].notna().any():
        tren = d.groupby("BULAN").size()
        tren.index = tren.index.astype(int)
        tren = tren.sort_index()

    # rekap harian + pola hari dalam pekan
    daily = None
    dow_avg = None
    if total and d["TGL PENGIRIMAN"].notna().any():
        _t = d.dropna(subset=["TGL PENGIRIMAN"])
        daily = _t.groupby(_t["TGL PENGIRIMAN"].dt.normalize()).size().sort_index()
        # rata-rata transaksi untuk tiap hari dalam pekan (Senin..Minggu)
        _per_day = _t.groupby([_t["TGL PENGIRIMAN"].dt.dayofweek,
                               _t["TGL PENGIRIMAN"].dt.normalize()]).size()
        dow_avg = _per_day.groupby(level=0).mean().sort_index()

    # highlight otomatis
    hl = []
    if total:
        hl.append(f"Tingkat penyelesaian {dec(done/total*100)}% dari {nf(total)} transaksi.")
        hl.append(f"Pembatalan {dec(cancel/total*100)}% — {nf(cancel)} unit tidak jadi dikerjakan.")
        if pending:
            top_pc = d[d["STATUS_BUCKET"] == "PENDING"]["CABANG"].value_counts()
            if len(top_pc):
                hl.append(f"{nf(pending)} unit masih tertahan; terbanyak di {top_pc.index[0]} "
                          f"({int(top_pc.iloc[0])} unit).")
        else:
            hl.append("Tidak ada unit yang tertahan pada periode ini.")
        if tren is not None and len(tren) > 1:
            bmax = tren.idxmax()
            hl.append(f"Bulan tersibuk: {BULAN_NAMES[bmax]} dengan {nf(tren.max())} transaksi.")

    footer_txt = (f"Sumber: {nf(total_raw_rows)} baris mentah → {nf(total_unique_all)} transaksi unik "
                  f"(baris identik dihitung satu). Dibuat otomatis dari dashboard "
                  f"{date.today().strftime('%d/%m/%Y')}.")

    # kesimpulan otomatis
    kes = []
    if total:
        kes.append(("Jaga tingkat penyelesaian",
                    f"Saat ini {dec(done/total*100)}% transaksi tuntas ({nf(done)} unit). "
                    f"Pertahankan kapasitas dan pastikan tidak turun saat volume naik."))
        cpct = cancel / total * 100
        c_sub = d[d["STATUS_BUCKET"] == "CANCEL"]
        c_ker = c_sub["KERUSAKAN"].value_counts()
        if len(c_ker):
            kes.append(("Telusuri penyebab pembatalan",
                        f"Pembatalan {dec(cpct)}% ({nf(cancel)} unit), paling sering pada kerusakan "
                        f"{c_ker.index[0]} ({nf(c_ker.iloc[0])} unit). Estimasi biaya lebih awal "
                        f"berpotensi menekan angka ini."))
        if pending:
            p_sub = d[d["STATUS_BUCKET"] == "PENDING"]
            p_cab = p_sub["CABANG"].value_counts()
            p_ker = p_sub["KERUSAKAN"].value_counts()
            kes.append(("Selesaikan unit yang tertahan",
                        f"{nf(pending)} unit pending. Prioritaskan {p_cab.index[0]} "
                        f"({nf(p_cab.iloc[0])} unit) dan kerusakan {p_ker.index[0]} "
                        f"({nf(p_ker.iloc[0])} unit) sebelum berkembang jadi komplain."))
        rek = (d.groupby("CABANG")["STATUS_BUCKET"]
               .apply(lambda x: (x == "CANCEL").sum() / len(x) * 100))
        rek = rek[d.groupby("CABANG").size() >= 30].sort_values(ascending=False)
        if len(rek):
            kes.append(("Tinjau cabang dengan pembatalan tertinggi",
                        f"{rek.index[0]} mencatat pembatalan {dec(rek.iloc[0])}%, "
                        f"di atas rata-rata {dec(cpct)}%. Perlu ditelusuri apakah soal harga, "
                        f"waktu tunggu, atau komunikasi."))

    meta = {
        "total": total, "done": done, "cancel": cancel, "pending": pending,
        "lainnya": lainnya, "avg_day": avg_day, "period_days": pdays,
        "tren": tren, "highlights": hl, "footer": footer_txt,
        "subjudul": subjudul, "kesimpulan": kes,
        "daily": daily, "dow_avg": dow_avg,
    }

    # --- rangkai slide ---
    _slide_cover(prs, periode, cabang_txt, total_unique_all, penyusun)
    if "ringkasan" in sertakan:
        _slide_ringkasan(prs, d, meta)
    if "status" in sertakan:
        _slide_status(prs, d, meta)
    if "banding_bulan" in sertakan:
        _sf_bd = sales_filtered
        if _sf_bd is not None and not _sf_bd.empty:
            _sf_bd = _sf_bd.copy()
            if "TGL" not in _sf_bd.columns:
                _sf_bd["TGL"] = pd.to_datetime(_sf_bd["TGL FAKTUR"], errors="coerce")
            if "MODAL" not in _sf_bd.columns:
                _sf_bd["MODAL"] = pd.to_numeric(_sf_bd.get("HARGA BELI", 0),
                                                errors="coerce").fillna(0)
            if "LABA" not in _sf_bd.columns:
                _sf_bd["LABA"] = _sf_bd["TOTAL HARGA"] - _sf_bd["MODAL"]
        _slide_banding_bulan(prs, d, _sf_bd, meta)
    if "harian" in sertakan:
        _slide_harian(prs, d, meta)
    if "hari_tertinggi" in sertakan:
        _slide_hari_tertinggi(prs, d, meta)
    if "cabang" in sertakan and f_cabang == "Semua Cabang":
        _slide_cabang(prs, d, meta)
    if "pending" in sertakan:
        _slide_status_detail(prs, d[d["STATUS_BUCKET"] == "PENDING"], meta,
                             judul="Unit Tertahan (Pending)", warna=AMBER, palette=RED,
                             jenis_func=_jenis_pending, jenis_label="Pending",
                             catatan="Pending adalah beban kerja yang harus segera diurai — "
                                     "makin lama tertahan, makin tinggi risiko komplain customer.")
    if "done" in sertakan:
        _slide_status_detail(prs, d[d["STATUS_BUCKET"] == "DONE"], meta,
                             judul="Penyelesaian (Done)", warna=GREEN_B, palette=GREEN,
                             jenis_func=lambda s: "Selesai", jenis_label="Done",
                             catatan=None)
    if "cancel" in sertakan:
        _slide_status_detail(prs, d[d["STATUS_BUCKET"] == "CANCEL"], meta,
                             judul="Pembatalan (Cancel)", warna=RED, palette=RGBColor(0x8B, 0x1E, 0x1E),
                             jenis_func=_jenis_cancel, jenis_label="Cancel",
                             catatan="Setiap pembatalan berpotensi menandakan persoalan di harga, "
                                     "waktu tunggu, ketersediaan sparepart, atau komunikasi.")
    # ---- slide berbasis data penjualan (opsional) ----
    sf = sales_filtered
    if sf is not None and not sf.empty:
        sf = sf.copy()
        if 'MODAL' not in sf.columns:
            sf['MODAL'] = pd.to_numeric(sf.get('HARGA BELI', 0), errors='coerce').fillna(0)
        if 'LABA' not in sf.columns:
            sf['LABA'] = sf['TOTAL HARGA'] - sf['MODAL']
        if 'TGL' not in sf.columns:
            sf['TGL'] = pd.to_datetime(sf['TGL FAKTUR'], errors='coerce')

        if "penjualan" in sertakan:
            _slide_penjualan(prs, sf, meta)

        if "mlf" in sertakan:
            nb = sf['NAMA BARANG'].astype(str).str.upper()
            _slide_mlf(prs, sf[nb.str.contains('MLF', na=False)], meta)

        if "bagihasil" in sertakan:
            tm = dict(tarif_map or {})
            for k, v in TARIF_AWAL_PPT.items():
                tm.setdefault(k, v)
            tm.setdefault('Lainnya', 30.0)

            kat = sf['KATEGORI BARANG'].astype(str).str.strip().str.upper()
            jasa = sf[kat == 'JASA'].copy()

            # periode cutoff 24 -> 23 bila bulan penggajian ditentukan
            if periode_bagihasil and not jasa.empty:
                ta, bg = periode_bagihasil
                m_akhir, th_akhir = bg - 1, ta
                if m_akhir < 1:
                    m_akhir += 12
                    th_akhir -= 1
                m_awal, th_awal = m_akhir - 1, th_akhir
                if m_awal < 1:
                    m_awal += 12
                    th_awal -= 1
                a = pd.Timestamp(th_awal, m_awal, 24)
                b = pd.Timestamp(th_akhir, m_akhir, 23)
                jasa = jasa[(jasa['TGL'] >= a) & (jasa['TGL'] <= b)]
                periode_txt = (f"Gaji {BULAN_NAMES[bg]} {ta} "
                               f"({a.day} {BULAN_NAMES[a.month]} – "
                               f"{b.day} {BULAN_NAMES[b.month]} {b.year})")
            else:
                periode_txt = periode

            if not jasa.empty:
                if 'TEKNISI' not in jasa.columns:
                    fin = jasa.get('NAMA TEKNISI (FINAL)')
                    asli = jasa.get('NAMA TEKNISI')
                    tek = (fin.fillna(asli) if fin is not None else asli)
                    jasa['TEKNISI'] = (tek.astype(str).str.strip().str.upper()
                                       .replace({'NAN': '', 'NONE': ''}))
                    jasa.loc[jasa['TEKNISI'] == '', 'TEKNISI'] = 'TIDAK ADA TEKNISI'
                jasa['TARIF_LABEL'] = jasa['NAMA BARANG'].map(
                    lambda s: _label_tarif(s, str(prioritas).upper()))
                jasa['TARIF'] = jasa['TARIF_LABEL'].map(
                    lambda k: tm.get(k, 0.0) / 100.0).fillna(0.0)
                jasa['BAGI_HASIL'] = jasa['TOTAL HARGA'] * jasa['TARIF']
                jasa['FLAT'] = jasa['TOTAL HARGA'] * (tarif_flat / 100.0)

            _slide_bagihasil(prs, jasa, meta, tarif_map=tm, tarif_flat=tarif_flat,
                             prioritas=prioritas, periode_txt=periode_txt)

    if "penutup" in sertakan and kes:
        _slide_penutup(prs, meta)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
